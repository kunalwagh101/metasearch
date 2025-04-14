# metasearch/extractors.py

import os
import json
import subprocess
import tarfile
import zipfile
from pathlib import Path
from datetime import datetime


try:
    import chardet
except ImportError:
    chardet = None


_EXTRACTOR_REGISTRY = {}

def register_extractor(file_extension, extractor_function):
    """
    Register a custom extractor for a specific file extension.
    """
    _EXTRACTOR_REGISTRY[file_extension.lower()] = extractor_function

def get_extractor_for(file_path):
    """
    Return the extractor for the file based on its extension;
    if none is registered, use the generic extractor.
    """
    extension = Path(file_path).suffix.lower()
    return _EXTRACTOR_REGISTRY.get(extension, extract_generic_metadata)

def inherent_metadata(file_path):
    """
    Extract inherent metadata (using pathlib.Path.stat()).
    Always returns a dictionary that includes "file_path" and, if possible,
    "file_name", "size_bytes", "created", "modified", and "access_time".
    """
    metadata = {"file_path": file_path}
    try:
        p = Path(file_path)
        stat_info = p.stat()
        metadata["file_path"] = str(p.resolve())
        metadata["file_name"] = p.name
        metadata["size_bytes"] = stat_info.st_size
        metadata["created"] = datetime.fromtimestamp(stat_info.st_ctime).isoformat()
        metadata["modified"] = datetime.fromtimestamp(stat_info.st_mtime).isoformat()
        metadata["access_time"] = datetime.fromtimestamp(stat_info.st_atime).isoformat()
        metadata["owner_uid"] = stat_info.st_uid
        metadata["group_gid"] = stat_info.st_gid
        metadata["permissions"] = oct(stat_info.st_mode)
    except Exception as e:
        metadata["inherent_error"] = str(e)
    return metadata


def extract_image_metadata(file_path):
    metadata = inherent_metadata(file_path)
    metadata["file_type"] = "image"
    try:
        from PIL import Image
        img = Image.open(file_path)
        metadata["width"], metadata["height"] = img.size
        metadata["mode"] = img.mode
    except Exception as e:
        metadata["image_error"] = f"Pillow error: {e}"
    try:
        with open(file_path, "rb") as f:
            import exifread
            exif = exifread.process_file(f, details=False)
            metadata["exif"] = {k: str(v) for k, v in exif.items()}
    except Exception as e:
        metadata["exif_error"] = f"ExifRead error: {e}"
    return metadata

def extract_audio_metadata(file_path):
    metadata = inherent_metadata(file_path)
    metadata["file_type"] = "audio"
    try:
        from mutagen import File as MutagenFile
        audio = MutagenFile(file_path)
        if audio is not None:
            if audio.tags:
                for tag in audio.tags.keys():
                    metadata[tag] = str(audio.tags.get(tag))
            if hasattr(audio.info, "length"):
                metadata["duration"] = audio.info.length
            if hasattr(audio.info, "bitrate"):
                metadata["bitrate"] = audio.info.bitrate
        else:
            metadata["audio_error"] = "Unsupported audio format or missing tags."
    except Exception as e:
        metadata["audio_error"] = str(e)
    return metadata

def extract_video_metadata(file_path):
    metadata = inherent_metadata(file_path)
    metadata["file_type"] = "video"
    try:
        cmd = [
            "ffprobe",
            "-v", "quiet",
            "-print_format", "json",
            "-show_format",
            "-show_streams",
            file_path
        ]
        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        if result.returncode == 0:
            info = json.loads(result.stdout)
            metadata["ffprobe"] = info
        else:
            metadata["video_error"] = result.stderr
    except Exception as e:
        metadata["video_error"] = str(e)
    return metadata

def extract_pdf_metadata(file_path):
    metadata = inherent_metadata(file_path)
    metadata["file_type"] = "pdf"
    try:
        import fitz
        doc = fitz.open(file_path)
        metadata["page_count"] = doc.page_count
        if doc.page_count > 0:
            page = doc[0]
            metadata["text_snippet"] = page.get_text()
        metadata.update(doc.metadata)
        doc.close()
    except Exception as e:
        metadata["pdf_error"] = str(e)
    return metadata

def extract_docx_metadata(file_path):
    metadata = inherent_metadata(file_path)
    metadata["file_type"] = "docx"
    try:
        import docx
        document = docx.Document(file_path)
        full_text = [para.text for para in document.paragraphs if para.text]
        metadata["full_text"] = "\n".join(full_text)
        cp = document.core_properties
        metadata["author"] = cp.author
        metadata["title"] = cp.title
        metadata["subject"] = cp.subject
        metadata["keywords"] = cp.keywords
        metadata["comments"] = cp.comments
    except Exception as e:
        metadata["docx_error"] = str(e)
    return metadata

def extract_xlsx_metadata(file_path):
    metadata = inherent_metadata(file_path)
    metadata["file_type"] = "xlsx"
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True)
        metadata["sheet_count"] = len(wb.sheetnames)
        metadata["sheets"] = wb.sheetnames
        wb.close()
    except Exception as e:
        metadata["xlsx_error"] = str(e)
    return metadata

def extract_pptx_metadata(file_path):
    metadata = inherent_metadata(file_path)
    metadata["file_type"] = "pptx"
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        metadata["slide_count"] = len(prs.slides)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    full_text.append(shape.text)
        metadata["full_text"] = "\n".join(full_text)
    except Exception as e:
        metadata["pptx_error"] = str(e)
    return metadata

def extract_text_metadata(file_path):
    metadata = inherent_metadata(file_path)
    metadata["file_type"] = "text"
    try:
        with open(file_path, "rb") as f:
            raw_data = f.read()
        try:
            text = raw_data.decode("utf-8")
        except UnicodeDecodeError:
            if chardet:
                detected = chardet.detect(raw_data)
                encoding = detected.get("encoding", "utf-8")
                text = raw_data.decode(encoding, errors="replace")
            else:
                text = raw_data.decode("utf-8", errors="replace")
        metadata["full_text"] = text
    except Exception as e:
        metadata["text_error"] = str(e)
    return metadata

def extract_archive_metadata(file_path):
    metadata = inherent_metadata(file_path)
    metadata["file_type"] = "archive"
    archive_list = []
    try:
        if file_path.lower().endswith(".zip"):
            import zipfile
            with zipfile.ZipFile(file_path, 'r') as archive:
                archive_list = archive.namelist()
        elif file_path.lower().endswith((".tar", ".gz", ".tgz")):
            with tarfile.open(file_path, 'r') as archive:
                archive_list = archive.getnames()
        else:
            archive_list = ["Archive type not supported"]
    except Exception as e:
        metadata["archive_error"] = str(e)
    metadata["contained_files"] = archive_list
    return metadata

def extract_generic_metadata(file_path):
    metadata = inherent_metadata(file_path)
    metadata["file_type"] = "binary"
    return metadata

# Register extractors for common file types.
for ext in [".jpg", ".jpeg", ".png", ".gif", ".tiff"]:
    register_extractor(ext, extract_image_metadata)
for ext in [".mp3", ".flac", ".ogg", ".m4a"]:
    register_extractor(ext, extract_audio_metadata)
for ext in [".mp4", ".avi", ".mkv", ".mov"]:
    register_extractor(ext, extract_video_metadata)
register_extractor(".pdf", extract_pdf_metadata)
for ext in [".docx"]:
    register_extractor(ext, extract_docx_metadata)
for ext in [".xlsx"]:
    register_extractor(ext, extract_xlsx_metadata)
for ext in [".pptx"]:
    register_extractor(ext, extract_pptx_metadata)
for ext in [".txt", ".py", ".java", ".c", ".md", ".json", ".xml"]:
    register_extractor(ext, extract_text_metadata)
for ext in [".zip", ".tar", ".gz", ".tgz"]:
    register_extractor(ext, extract_archive_metadata)
